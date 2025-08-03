# أزل استيراد transformers و torch
from nltk.tokenize import sent_tokenize
import random
import fitz  # PyMuPDF
import pptx

# تحميل نموذج BioBERT
tokenizer = AutoTokenizer.from_pretrained("dmis-lab/biobert-v1.1")
model = AutoModelForQuestionAnswering.from_pretrained("dmis-lab/biobert-v1.1")

def extract_text_from_file(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(file_path)
    return ""

def extract_text_from_pdf(file_path):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        text = f"Error reading PDF: {e}"
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Error reading PPTX: {e}"
    return text

def calculate_difficulty(text):
    medical_terms = ["diagnosis", "treatment", "symptoms", "clinical", "patient", "disease"]
    term_count = sum(1 for word in text.lower().split() if word in medical_terms)
    return "easy" if term_count < 3 else "medium" if term_count < 6 else "hard"

def generate_mcq(text, num_questions):
    sentences = [s for s in sent_tokenize(text) if len(s.split()) > 5]
    random.shuffle(sentences)
    questions = []
    
    for sentence in sentences[:num_questions]:
        difficulty = calculate_difficulty(sentence)
        inputs = tokenizer(
            f"Generate a {difficulty} medical question about: {sentence}",
            return_tensors="pt",
            max_length=512,
            truncation=True
        )
        with torch.no_grad():
            outputs = model(**inputs)
        question = tokenizer.decode(outputs.start_logits.argmax(), skip_special_tokens=True)
        
        options = [
            "Correct answer (placeholder)",
            "Incorrect option 1",
            "Incorrect option 2",
            "Incorrect option 3"
        ]
        random.shuffle(options)
        correct = options.index("Correct answer (placeholder)")
        
        questions.append({
            "question": f"{question} (الصعوبة: {difficulty})",
            "options": options,
            "correct": correct
        })
    return questions
